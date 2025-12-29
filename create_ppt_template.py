#!/usr/bin/env python3
"""
Create VLM-Metrics PowerPoint Template - Professional Edition
"""

from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Professional color palette
DARK_NAVY = RGBColor(25, 42, 86)  # #192A56 - Sophisticated navy
CHARCOAL = RGBColor(44, 62, 80)  # #2C3E50 - Professional gray-blue
SLATE_BLUE = RGBColor(52, 73, 94)  # #34495E - Muted slate
ACCENT_BLUE = RGBColor(41, 128, 185)  # #2980B9 - Professional blue accent
LIGHT_GRAY = RGBColor(236, 240, 241)  # #ECF0F1 - Subtle background
MID_GRAY = RGBColor(149, 165, 166)  # #95A5A6 - Text gray
WHITE = RGBColor(255, 255, 255)
TEXT_DARK = RGBColor(33, 47, 61)  # #212F3D - Professional text

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)  # 16:9 aspect ratio

def add_title_slide(prs):
    """Add professional title slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Clean white background
    bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = WHITE
    bg_shape.line.fill.background()
    
    # Subtle left accent bar
    accent_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(5.625))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_BLUE
    accent_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.5), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Presentation Title"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_NAVY
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.1), Inches(8.5), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Subtitle or presentation context"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = CHARCOAL
    
    # Bottom branding section
    brand_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.9), Inches(4), Inches(0.5))
    brand_frame = brand_box.text_frame
    brand_frame.text = "VLM-Metrics | Responsible AI Governance"
    brand_para = brand_frame.paragraphs[0]
    brand_para.font.size = Pt(11)
    brand_para.font.color.rgb = MID_GRAY
    
    # Date/presenter
    presenter_box = slide.shapes.add_textbox(Inches(6), Inches(4.9), Inches(3.2), Inches(0.5))
    presenter_frame = presenter_box.text_frame
    presenter_frame.text = "Presenter Name | Date"
    presenter_para = presenter_frame.paragraphs[0]
    presenter_para.font.size = Pt(11)
    presenter_para.font.color.rgb = MID_GRAY
    presenter_para.alignment = PP_ALIGN.RIGHT
    
    return slide

def add_content_slide(prs):
    """Add professional content slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = WHITE
    bg_shape.line.fill.background()
    
    # Minimal top accent line
    accent_line = slide.shapes.add_shape(1, Inches(0), Inches(0.5), Inches(10), Inches(0.02))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = ACCENT_BLUE
    accent_line.line.fill.background()
    
    # Header with branding
    brand_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(3), Inches(0.3))
    brand_frame = brand_box.text_frame
    brand_frame.text = "VLM-Metrics"
    brand_para = brand_frame.paragraphs[0]
    brand_para.font.size = Pt(9)
    brand_para.font.color.rgb = MID_GRAY
    
    # Slide number
    number_box = slide.shapes.add_textbox(Inches(9.2), Inches(0.15), Inches(0.3), Inches(0.3))
    number_frame = number_box.text_frame
    number_frame.text = "1"
    number_para = number_frame.paragraphs[0]
    number_para.font.size = Pt(9)
    number_para.font.color.rgb = MID_GRAY
    number_para.alignment = PP_ALIGN.RIGHT
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Slide Title"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_NAVY
    
    # Content area
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(3.5))
    content_frame = content_box.text_frame
    content_frame.text = "• First key point or insight\n• Second key point or insight\n• Third key point or insight\n• Fourth key point or insight"
    for para in content_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = TEXT_DARK
        para.space_after = Pt(14)
        para.line_spacing = 1.3
    
    return slide

def add_section_slide(prs):
    """Add professional section divider"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Light gray background
    bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = LIGHT_GRAY
    bg_shape.line.fill.background()
    
    # Left accent panel
    accent_panel = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.4), Inches(5.625))
    accent_panel.fill.solid()
    accent_panel.fill.fore_color.rgb = ACCENT_BLUE
    accent_panel.line.fill.background()
    
    # Section number (optional)
    section_num_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(0.4))
    section_num_frame = section_num_box.text_frame
    section_num_frame.text = "01"
    section_num_para = section_num_frame.paragraphs[0]
    section_num_para.font.size = Pt(72)
    section_num_para.font.bold = True
    section_num_para.font.color.rgb = RGBColor(220, 220, 220)
    
    # Section name
    section_box = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(8), Inches(0.8))
    section_frame = section_box.text_frame
    section_frame.text = "Section Name"
    section_para = section_frame.paragraphs[0]
    section_para.font.size = Pt(40)
    section_para.font.bold = True
    section_para.font.color.rgb = DARK_NAVY
    
    # VLM-Metrics branding
    brand_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(3), Inches(0.3))
    brand_frame = brand_box.text_frame
    brand_frame.text = "VLM-Metrics"
    brand_para = brand_frame.paragraphs[0]
    brand_para.font.size = Pt(9)
    brand_para.font.color.rgb = MID_GRAY
    
    return slide

def add_two_column_slide(prs):
    """Add professional two-column layout"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = WHITE
    bg_shape.line.fill.background()
    
    # Top accent line
    accent_line = slide.shapes.add_shape(1, Inches(0), Inches(0.5), Inches(10), Inches(0.02))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = ACCENT_BLUE
    accent_line.line.fill.background()
    
    # Header
    brand_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(3), Inches(0.3))
    brand_frame = brand_box.text_frame
    brand_frame.text = "VLM-Metrics"
    brand_para = brand_frame.paragraphs[0]
    brand_para.font.size = Pt(9)
    brand_para.font.color.rgb = MID_GRAY
    
    number_box = slide.shapes.add_textbox(Inches(9.2), Inches(0.15), Inches(0.3), Inches(0.3))
    number_frame = number_box.text_frame
    number_frame.text = "2"
    number_para = number_frame.paragraphs[0]
    number_para.font.size = Pt(9)
    number_para.font.color.rgb = MID_GRAY
    number_para.alignment = PP_ALIGN.RIGHT
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Slide Title"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_NAVY
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.5), Inches(3.5))
    left_frame = left_box.text_frame
    left_frame.text = "• Key insight one\n• Key insight two\n• Key insight three\n• Key insight four"
    for para in left_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = TEXT_DARK
        para.space_after = Pt(12)
        para.line_spacing = 1.3
    
    # Right column (chart/image placeholder)
    right_box = slide.shapes.add_shape(1, Inches(5.2), Inches(1.6), Inches(4.3), Inches(3.5))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = LIGHT_GRAY
    right_box.line.color.rgb = MID_GRAY
    right_box.line.width = Pt(1)
    
    placeholder_text = right_box.text_frame
    placeholder_text.text = "[Chart or Image]"
    placeholder_para = placeholder_text.paragraphs[0]
    placeholder_para.font.size = Pt(16)
    placeholder_para.font.color.rgb = MID_GRAY
    placeholder_para.alignment = PP_ALIGN.CENTER
    placeholder_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return slide

def add_closing_slide(prs):
    """Add professional closing slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = WHITE
    bg_shape.line.fill.background()
    
    # Left accent bar
    accent_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(5.625))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_BLUE
    accent_bar.line.fill.background()
    
    # Thank you text
    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(0.8))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "Thank You"
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.font.size = Pt(48)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = DARK_NAVY
    thanks_para.alignment = PP_ALIGN.CENTER
    
    # Questions prompt
    questions_box = slide.shapes.add_textbox(Inches(1), Inches(2.9), Inches(8), Inches(0.5))
    questions_frame = questions_box.text_frame
    questions_frame.text = "Questions & Discussion"
    questions_para = questions_frame.paragraphs[0]
    questions_para.font.size = Pt(22)
    questions_para.font.color.rgb = CHARCOAL
    questions_para.alignment = PP_ALIGN.CENTER
    
    # Contact information
    contact_box = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(0.8))
    contact_frame = contact_box.text_frame
    contact_frame.text = "contact@vlm-metrics.com\nwww.vlm-metrics.com"
    for para in contact_frame.paragraphs:
        para.font.size = Pt(14)
        para.font.color.rgb = MID_GRAY
        para.alignment = PP_ALIGN.CENTER
        para.space_after = Pt(6)
    
    # Bottom branding
    brand_box = slide.shapes.add_textbox(Inches(2.5), Inches(5), Inches(5), Inches(0.4))
    brand_frame = brand_box.text_frame
    brand_frame.text = "VLM-Metrics | Responsible AI Governance"
    brand_para = brand_frame.paragraphs[0]
    brand_para.font.size = Pt(11)
    brand_para.font.color.rgb = MID_GRAY
    brand_para.alignment = PP_ALIGN.CENTER
    
    return slide

# Create all template slides
print("Creating VLM-Metrics Professional PowerPoint Template...")
add_title_slide(prs)
add_content_slide(prs)
add_section_slide(prs)
add_two_column_slide(prs)
add_closing_slide(prs)

# Save the presentation
timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M")
output_file = f"business_presentations/VLM-Metrics_Template_{timestamp}.pptx"
prs.save(output_file)
print(f"✓ Professional template created: {output_file}")
print(f"  - 5 refined slide templates")
print(f"  - Title Slide (clean, minimal)")
print(f"  - Content Slide (professional layout)")
print(f"  - Section Divider (sophisticated)")
print(f"  - Two-Column Layout (balanced)")
print(f"  - Closing Slide (elegant)")
print(f"\nDesign: Muted navy palette, corporate typography, minimal accents")
